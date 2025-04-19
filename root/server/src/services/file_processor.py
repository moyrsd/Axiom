import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from pandas import read_csv, read_excel, read_json
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
from typing import List
from langchain.text_splitter import RecursiveCharacterTextSplitter
from ..services import vector_store



def process_file(file_path: str) -> str:
    text = ""
    ext = os.path.splitext(file_path)[-1].lower()
    
    if ext == ".pdf":
        text += process_pdf(file_path)
    elif ext == ".docx":
        doc = Document(file_path)
        text += "\n".join([para.text for para in doc.paragraphs])
    elif ext == ".pptx":
        prs = Presentation(file_path)
        text += "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif ext in (".xlsx", ".xls"):
        df = read_excel(file_path)
        text += df.to_string()
    elif ext == ".csv":
        df = read_csv(file_path)
        text += df.to_string()
    elif ext == ".json":
        with open(file_path) as f:
            data = read_json(f)
            text += str(data)
    elif ext in (".png", ".jpg", ".jpeg"):
        text += pytesseract.image_to_string(Image.open(file_path))
    elif ext == ".txt":
        with open(file_path) as f:
            text += f.read()
    
    return text

def process_pdf(file_path: str) -> str:
    text = ""
    try:  # First try text extraction
        with open(file_path, "rb") as file:
            pdf_reader = PdfReader(file)
            text += "".join([page.extract_text() or "" for page in pdf_reader.pages])
    except:  # Fallback to OCR
        images = convert_from_path(file_path)
        text += "".join([pytesseract.image_to_string(img) for img in images])
    return text

def get_text_chunks(text: str) -> List[str]:
    """Split text into manageable chunks"""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    return text_splitter.split_text(text)





def process_files(temp_paths: List[str]):
    try:
        raw_text = ""
        for path in temp_paths:
            raw_text += process_file(path) + "\n\n"
        
        text_chunks = get_text_chunks(raw_text)
        vector_store.create_vector_store(text_chunks)
        
        # Cleanup
        for path in temp_paths:
            os.remove(path)
            
    except Exception as e:
        print(f"Processing error: {str(e)}")




