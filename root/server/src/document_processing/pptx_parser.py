from pptx import Presentation
from pdf2image import convert_from_path
import pytesseract


def process_pptx(file_path: str) -> str:
    text = ""
    try:  # First try text extraction
        with open(file_path, "rb") as file:
            prs =  Presentation(file)
            text += "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    except:  # Fallback to OCR
        images = convert_from_path(file_path)
        text += "".join([pytesseract.image_to_string(img) for img in images])
    return text