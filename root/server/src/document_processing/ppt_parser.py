from pptx import Presentation
import os
from io import BytesIO
from PIL import Image
from ..services import llm_calls
from ..prompts import ocr_prompt

def process_pptx(file_path: str) -> list:
    """Extracts text + OCR from images in PowerPoint slides."""
    extracted_content = []
    
    try:
        prs = Presentation(file_path)
        filename = os.path.basename(file_path)
        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                if shape.shape_type == 13:  
                    try:
                        image = shape.image
                        img_bytes = image.blob
    
                        img = Image.open(BytesIO(img_bytes))
                        prompt = ocr_prompt.prompt_ocr
                        ocr_client = llm_calls.LlmCalls()
                        ocr_text = ocr_client.llm_ocr(img,prompt)
                        if ocr_text.strip():
                            slide_text.append(f"[IMAGE OCR]:\n{ocr_text.strip()}")
                            
                    except Exception as e:
                        print(f"OCR failed on slide {slide_number}: {str(e)}")

            extracted_content.append({
                "page_content": "\n".join(slide_text),
                "meta_data": {
                    "source": f"{filename}__pageno-{slide_number}"
                }
            })
            
    except Exception as e:
        print(f"PPT processing error: {str(e)}")
    print(extracted_content)
    return extracted_content


# https://python-pptx.readthedocs.io/en/latest/